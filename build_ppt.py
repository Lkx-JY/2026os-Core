#!/usr/bin/env python3
"""
Build competition presentation PPT for:
Linux Kernel Crash → Root Cause → Patch Matching System
Based on RAG + Kernel Domain Knowledge + LLM
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import os

# ============================================================
# CONSTANTS & COLOR PALETTE
# ============================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
C_DARK_BG    = RGBColor(0x0F, 0x1B, 0x2D)  # Deep navy
C_BLUE       = RGBColor(0x1E, 0x3A, 0x5F)   # Secondary blue
C_ACCENT     = RGBColor(0xE8, 0x79, 0x2B)   # Orange accent
C_GREEN      = RGBColor(0x2E, 0x8B, 0x57)   # Success green
C_RED        = RGBColor(0xDC, 0x35, 0x35)   # Alert red
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG   = RGBColor(0xF7, 0xFA, 0xFC)
C_LIGHT_GRAY = RGBColor(0xED, 0xF2, 0xF7)
C_GRAY       = RGBColor(0x4A, 0x55, 0x68)
C_TEXT       = RGBColor(0x1A, 0x20, 0x2C)
C_GOLD       = RGBColor(0xD6, 0x9E, 0x2E)

# Fonts
F_TITLE = '微软雅黑'
F_BODY  = '微软雅黑'
F_EN    = 'Arial'

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout (index 6)
BLANK_LAYOUT = prs.slide_layouts[6]

# ============================================================
# HELPER FUNCTIONS
# ============================================================
def add_blank_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def add_textbox(slide, left, top, width, height, text="", font_name=F_BODY,
                font_size=Pt(14), bold=False, color=C_TEXT, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = alignment
    tf.paragraphs[0].line_spacing = line_spacing
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    # Set East Asian font
    for rPr in txBox._element.findall('.//' + qn('a:rPr')):
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font_name)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_font=F_BODY,
                          default_size=Pt(14), default_color=C_TEXT):
    """
    Add a text box with multiple formatted lines.
    lines: list of dicts: {'text': str, 'font_name': ..., 'font_size': ..., 'bold': ..., 'color': ..., 'alignment': ..., 'line_spacing': ...}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        para.alignment = line.get('alignment', PP_ALIGN.LEFT)
        para.line_spacing = line.get('line_spacing', 1.15)
        para.space_after = line.get('space_after', Pt(2))

        run = para.add_run()
        run.text = line.get('text', '')
        run.font.name = line.get('font_name', default_font)
        run.font.size = line.get('font_size', default_size)
        run.font.bold = line.get('bold', False)
        run.font.color.rgb = line.get('color', default_color)

        # Set East Asian font
        rPr_elements = para._p.findall(qn('a:r') + '/' + qn('a:rPr'))
        for rPr in rPr_elements:
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', run.font.name)

    return txBox

def add_circle(slide, left, top, size, fill_color=None, text="", font_size=Pt(14), font_color=C_WHITE):
    """Add a circle with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.name = F_BODY
        run.font.size = font_size
        run.font.bold = True
        run.font.color.rgb = font_color
    return shape

def add_arrow_right(slide, left, top, width, height, color=C_ACCENT):
    """Add a right arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, height, color=C_BLUE, width_pt=Pt(2)):
    """Add a horizontal line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_page_number(slide, num, total=11):
    """Add page number at bottom right."""
    add_textbox(slide, Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.35),
                f"{num} / {total}", font_size=Pt(9), color=C_GRAY,
                alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, section_num, title_cn, title_en=None):
    """Add a consistent section header bar at top."""
    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), C_ACCENT)
    # Section number circle
    add_circle(slide, Inches(0.6), Inches(0.35), Inches(0.55), C_ACCENT,
               f"{section_num:02d}", Pt(18))
    # Title
    add_textbox(slide, Inches(1.35), Inches(0.35), Inches(10), Inches(0.55),
                title_cn, font_name=F_TITLE, font_size=Pt(26), bold=True, color=C_DARK_BG)
    if title_en:
        add_textbox(slide, Inches(1.35), Inches(0.85), Inches(10), Inches(0.35),
                    title_en, font_name=F_EN, font_size=Pt(12), color=C_GRAY)
    # Bottom line
    add_line(slide, Inches(0.6), Inches(1.3), Inches(12.1), Pt(2), C_LIGHT_GRAY)

def add_card(slide, left, top, width, height, title, content_lines,
             title_color=C_BLUE, bg_color=C_WHITE, border_color=C_LIGHT_GRAY):
    """Add a card with title and content."""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, bg_color)
    # Left accent bar
    add_line(slide, left + Inches(0.02), top + Inches(0.1), Inches(0.05), height - Inches(0.2), title_color)

    # Title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.35),
                title, font_name=F_TITLE, font_size=Pt(13), bold=True, color=title_color)

    # Content
    y_offset = top + Inches(0.5)
    for line in content_lines:
        add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.28),
                    line, font_size=Pt(10.5), color=C_GRAY)
        y_offset += Inches(0.25)

def add_table(slide, left, top, col_widths, headers, rows, header_bg=C_BLUE):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(w for w in col_widths)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.4) * n_rows)
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = header
        run.font.name = F_BODY
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = C_WHITE
        # Cell fill
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = para.add_run()
            run.text = str(val)
            run.font.name = F_BODY
            run.font.size = Pt(10)
            run.font.color.rgb = C_TEXT
            # Alternating row color
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT_BG

    return table_shape

# ============================================================
# SLIDE 1: COVER
# ============================================================
print("Building Slide 1 - Cover...")
slide = add_blank_slide()

# Background - dark gradient effect using two rectangles
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_DARK_BG)

# Decorative accent shapes
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), C_ACCENT)
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), C_ACCENT)

# Left decorative circle (large, subtle)
c1 = add_circle(slide, Inches(-1.5), Inches(3.5), Inches(4), RGBColor(0x1A, 0x2E, 0x4A))
c1.fill.solid()
c1.fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x4A)
add_textbox(slide, Inches(0), Inches(0), Inches(1), Inches(1), "", font_size=Pt(10))  # placeholder

# Right accent stripe
add_rect(slide, Inches(12.6), Inches(1.5), Inches(0.08), Inches(4.5), C_ACCENT)

# Linux penguin icon area (text-based)
add_textbox(slide, Inches(1.2), Inches(2.0), Inches(2.5), Inches(2.0),
            "🐧", font_size=Pt(72), alignment=PP_ALIGN.CENTER, color=C_WHITE)

# Main title
add_textbox(slide, Inches(1.2), Inches(3.8), Inches(11.0), Inches(1.0),
            "Linux 内核宕机自动诊断与补丁匹配系统",
            font_name=F_TITLE, font_size=Pt(40), bold=True, color=C_WHITE)

# Subtitle
add_textbox(slide, Inches(1.2), Inches(4.7), Inches(11.0), Inches(0.6),
            "基于 RAG + 内核领域知识 + LLM 的自动化补丁匹配方案",
            font_name=F_TITLE, font_size=Pt(20), color=C_ACCENT)

# Key value proposition
add_textbox(slide, Inches(1.2), Inches(5.5), Inches(10.5), Inches(0.5),
            '核心突破：解决「百万级 Commit 中快速精准匹配宕机修复补丁」难题',
            font_name=F_BODY, font_size=Pt(14), color=RGBColor(0xA0, 0xAE, 0xC0))

# Bottom info bar
add_multiline_textbox(slide, Inches(1.2), Inches(6.5), Inches(10), Inches(0.7), [
    {'text': '项目团队：XXX    |    指导教师：XXX    |    2026年6月',
     'font_size': Pt(12), 'color': C_GRAY},
])

# Right side: flow diagram (simple shapes)
# Crash icon area
add_rect(slide, Inches(9.5), Inches(1.8), Inches(3.0), Inches(1.5),
         RGBColor(0x1A, 0x2E, 0x4A), C_ACCENT, Pt(1))
add_textbox(slide, Inches(9.5), Inches(2.0), Inches(3.0), Inches(0.5),
            "💥 宕机日志", font_size=Pt(16), bold=True, color=C_RED,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.5), Inches(2.5), Inches(3.0), Inches(0.6),
            "dmesg / vmcore / 问题描述", font_size=Pt(10), color=C_GRAY,
            alignment=PP_ALIGN.CENTER)

# Arrow down
add_textbox(slide, Inches(10.7), Inches(3.5), Inches(0.6), Inches(0.8),
            "⬇", font_size=Pt(28), color=C_ACCENT, alignment=PP_ALIGN.CENTER)

# System box
add_rect(slide, Inches(9.5), Inches(4.2), Inches(3.0), Inches(1.5),
         RGBColor(0x1A, 0x2E, 0x4A), C_GREEN, Pt(1))
add_textbox(slide, Inches(9.5), Inches(4.4), Inches(3.0), Inches(0.5),
            "🔧 RAG 智能诊断引擎", font_size=Pt(16), bold=True, color=C_GREEN,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.5), Inches(4.9), Inches(3.0), Inches(0.6),
            "根因分析 → 向量检索 → 补丁匹配", font_size=Pt(10), color=C_GRAY,
            alignment=PP_ALIGN.CENTER)

# Arrow down
add_textbox(slide, Inches(10.7), Inches(5.85), Inches(0.6), Inches(0.8),
            "⬇", font_size=Pt(28), color=C_GREEN, alignment=PP_ALIGN.CENTER)

# Fix icon area
add_rect(slide, Inches(9.5), Inches(6.5), Inches(3.0), Inches(0.8),
         RGBColor(0x1A, 0x2E, 0x4A), C_GOLD, Pt(1))
add_textbox(slide, Inches(9.5), Inches(6.55), Inches(3.0), Inches(0.7),
            "✅ Top-N 精准补丁推荐", font_size=Pt(14), bold=True, color=C_GOLD,
            alignment=PP_ALIGN.CENTER)

add_page_number(slide, 1)

# ============================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================
print("Building Slide 2 - TOC...")
slide = add_blank_slide()

# Background
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), C_ACCENT)

# Title
add_textbox(slide, Inches(1.0), Inches(0.5), Inches(5), Inches(0.7),
            "汇报大纲", font_name=F_TITLE, font_size=Pt(36), bold=True, color=C_DARK_BG)
add_line(slide, Inches(1.0), Inches(1.2), Inches(3), Pt(3), C_ACCENT)

# TOC items with icons and numbers
toc_items = [
    ("01", "赛题背景与核心挑战", "🔴"),
    ("02", "系统总体架构（离线 + 在线）", "🏗️"),
    ("03", "核心创新：对称 Root Cause Embedding", "💡"),
    ("04", "离线阶段：Commit 知识库构建", "📦"),
    ("05", "在线阶段：四阶段检索与诊断", "🔍"),
    ("06", "技术亮点与性能指标", "📊"),
    ("07", "总结与应用价值", "🎯"),
]

for i, (num, title, icon) in enumerate(toc_items):
    y = Inches(1.8) + i * Inches(0.75)

    # Number circle
    add_circle(slide, Inches(1.2), y, Inches(0.45), C_BLUE, num, Pt(14))

    # Title
    add_textbox(slide, Inches(1.9), y + Inches(0.05), Inches(8), Inches(0.4),
                title, font_name=F_TITLE, font_size=Pt(18), bold=True, color=C_DARK_BG)

    # Dotted line
    add_line(slide, Inches(1.9), y + Inches(0.55), Inches(8.5), Pt(1), C_LIGHT_GRAY)

# Right side: decorative element
add_rect(slide, Inches(11.0), Inches(1.5), Inches(1.8), Inches(5.0), C_WHITE)
add_textbox(slide, Inches(11.2), Inches(2.0), Inches(1.4), Inches(4.0),
            "Linux\nKernel\nCrash\n→\nPatch\nMatch",
            font_name=F_EN, font_size=Pt(16), bold=True, color=C_ACCENT,
            alignment=PP_ALIGN.CENTER)

add_page_number(slide, 2)

# ============================================================
# SLIDE 3: BACKGROUND & CHALLENGES
# ============================================================
print("Building Slide 3 - Background & Challenges...")
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_section_header(slide, 1, "内核宕机诊断的行业痛点与核心挑战",
                   "Industry Pain Points & Core Challenges in Kernel Crash Diagnosis")

# Left column: Background pain points
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.4),
            "▎行业背景与痛点", font_size=Pt(16), bold=True, color=C_DARK_BG)

pain_points = [
    ("运维痛点", "145万+ Linux Kernel Commit，人工分析宕机需数小时/天，\n严重依赖资深内核专家，中小企业无法承担人力成本。"),
    ("语义鸿沟", '宕机日志呈现「运行时错误现象」（如 list_del 损坏），\n补丁描述「代码层修复逻辑」（如添加锁保护），两者语义不匹配。'),
    ("多模态输入", "生产环境宕机信息多样：vmcore 二进制内存镜像、\ndmesg 文本日志、运维人员自然语言描述，异构数据难统一。"),
    ("版本演化", "同一函数跨内核版本实现差异大（kernel 4.9 vs 6.13），\n补丁需精确匹配目标版本，避免引入新问题。"),
]

y = Inches(2.15)
for (label, desc) in pain_points:
    # Label badge
    badge = add_rounded_rect(slide, Inches(0.6), y, Inches(1.2), Inches(0.28), C_ACCENT)
    add_textbox(slide, Inches(0.65), y, Inches(1.1), Inches(0.28),
                label, font_size=Pt(9), bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, Inches(2.0), y - Inches(0.02), Inches(4.5), Inches(0.7),
                desc, font_size=Pt(10.5), color=C_GRAY)
    y += Inches(0.9)

# Right column: 4 Core Challenges (table)
add_textbox(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(0.4),
            "▎四大核心挑战", font_size=Pt(16), bold=True, color=C_DARK_BG)

challenge_headers = ["挑战维度", "核心问题", "量化指标/示例"]
challenge_cols = [Inches(1.6), Inches(2.5), Inches(1.8)]
challenge_rows = [
    ["语义非对称性", '宕机现象≠修复逻辑', '相似度仅0.3\n「list_del损坏」≠「加锁」'],
    ["海量检索效率", "百万级Commit秒级检索", "145万×28子系统\n×21 Bug类型"],
    ["版本演化敏感", "同函数跨版本实现差异", "kernel 4.9 vs 6.13\n函数逻辑不同"],
    ["多模态融合", "异构数据统一特征表示", "vmcore+dmesg\n+问题描述"],
]

add_table(slide, Inches(6.8), Inches(2.15), challenge_cols, challenge_headers, challenge_rows, C_RED)

# Bottom highlight
add_rounded_rect(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0), C_WHITE)
add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8),
            '💡 本项目核心目标：构建「输入宕机日志 → 自动根因分析 → 百万级Commit精准匹配 → 可解释诊断报告」的全自动智能诊断流水线',
            font_size=Pt(13), bold=True, color=C_BLUE)

add_page_number(slide, 3)

# ============================================================
# SLIDE 4: SYSTEM ARCHITECTURE
# ============================================================
print("Building Slide 4 - System Architecture...")
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_section_header(slide, 2, "宏观架构：离线治理 + 在线检索双阶段",
                   "System Architecture: Offline Indexing + Online Retrieval")

# Architecture diagram using shapes
# --- LEFT: Offline Phase (Blue) ---
offline_x = Inches(0.5)
offline_w = Inches(5.8)

add_rounded_rect(slide, offline_x, Inches(1.6), offline_w, Inches(5.3),
                 RGBColor(0xEB, 0xF4, 0xFF))  # Light blue bg

add_textbox(slide, offline_x + Inches(0.2), Inches(1.7), offline_w - Inches(0.4), Inches(0.4),
            "🔵 离线阶段：Commit 知识库构建", font_size=Pt(15), bold=True, color=C_BLUE)

offline_modules = [
    ("Linux Kernel\nGit Repo", "145万+ commits", C_BLUE),
    ("PyDriller\n流式遍历", "O(1)内存", C_BLUE),
    ("CommitRootCause\nBuilder", "3-5ms/条", C_BLUE),
    ("BGE-M3\n向量编码", "1024维", C_BLUE),
    ("Milvus/FAISS\n向量存储", "全量索引", C_BLUE),
]

off_x = offline_x + Inches(0.15)
for i, (name, desc, color) in enumerate(offline_modules):
    box = add_rounded_rect(slide, off_x, Inches(2.4), Inches(1.0), Inches(1.6), C_WHITE)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    add_textbox(slide, off_x, Inches(2.5), Inches(1.0), Inches(0.6),
                name, font_size=Pt(9), bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, off_x, Inches(3.2), Inches(1.0), Inches(0.4),
                desc, font_size=Pt(7.5), color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes
    if i < len(offline_modules) - 1:
        add_textbox(slide, off_x + Inches(1.0), Inches(3.0), Inches(0.15), Inches(0.3),
                    "→", font_size=Pt(16), bold=True, color=C_ACCENT, alignment=PP_ALIGN.CENTER)

    off_x += Inches(1.15)

# Key design points
add_textbox(slide, offline_x + Inches(0.2), Inches(4.3), offline_w - Inches(0.4), Inches(2.3),
            "关键设计：\n"
            "• 生成器模式流式处理，O(1)内存占用\n"
            "• 25条 DIFF_RULES 规则：锁/RCU/内存/引用计数等\n"
            "• 8条锁规则 + 4条引用计数规则 + 内核版本标注\n"
            "• 断点续跑支持，百万级全量索引可行",
            font_size=Pt(9.5), color=C_GRAY)

# --- RIGHT: Online Phase (Orange) ---
online_x = Inches(6.8)
online_w = Inches(6.0)

add_rounded_rect(slide, online_x, Inches(1.6), online_w, Inches(5.3),
                 RGBColor(0xFF, 0xF4, 0xE6))  # Light orange bg

add_textbox(slide, online_x + Inches(0.2), Inches(1.7), online_w - Inches(0.4), Inches(0.4),
            "🟠 在线阶段：宕机诊断与补丁匹配", font_size=Pt(15), bold=True, color=C_ACCENT)

online_modules = [
    ("dmesg/vmcore\n问题描述", "多模态输入", C_ACCENT),
    ("CrashFeature\n特征提取", "20+正则", C_ACCENT),
    ("RootCause\nAnalyzer", "28条规则", C_ACCENT),
    ("四阶段检索\nRecall→Judge", "100→5", C_ACCENT),
    ("补丁推荐\n可解释报告", "Top-3~5", C_ACCENT),
]

on_x = online_x + Inches(0.2)
for i, (name, desc, color) in enumerate(online_modules):
    box = add_rounded_rect(slide, on_x, Inches(2.4), Inches(1.05), Inches(1.6), C_WHITE)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    add_textbox(slide, on_x, Inches(2.5), Inches(1.05), Inches(0.6),
                name, font_size=Pt(9), bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, on_x, Inches(3.2), Inches(1.05), Inches(0.4),
                desc, font_size=Pt(7.5), color=C_GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(online_modules) - 1:
        add_textbox(slide, on_x + Inches(1.05), Inches(3.0), Inches(0.15), Inches(0.3),
                    "→", font_size=Pt(16), bold=True, color=C_GREEN, alignment=PP_ALIGN.CENTER)

    on_x += Inches(1.15)

# Online key points
add_textbox(slide, online_x + Inches(0.2), Inches(4.3), online_w - Inches(0.4), Inches(2.3),
            "关键设计：\n"
            "• 20+ Panic 模式正则 + LLM 深度特征提取\n"
            "• 28条专家规则 + 4层分层根因推断\n"
            "• 四阶段级联检索：Recall→Filter→Rerank→LLM Judge\n"
            "• 支持 fast(<100ms) / standard(<1s) / deep(2-10s) 模式",
            font_size=Pt(9.5), color=C_GRAY)

# Bottom tech stack bar
add_rounded_rect(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.32), C_BLUE)
add_textbox(slide, Inches(0.7), Inches(7.05), Inches(12), Inches(0.32),
            "技术栈：Python 3.12 | FastAPI | BGE-M3 | BGE-Reranker-v2 | Milvus/FAISS | PyDriller | drgn | Vue 3 | Docker | Ollama",
            font_size=Pt(9), bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 4)

# ============================================================
# SLIDE 5: CORE INNOVATION - Symmetric Root Cause Embedding
# ============================================================
print("Building Slide 5 - Core Innovation...")
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_section_header(slide, 3, "核心突破：解决「语义非对称」的对称嵌入方案",
                   "Core Innovation: Symmetric Root Cause Embedding")

# Problem side (Left) - Red tint
add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.8),
                 RGBColor(0xFF, 0xF5, 0xF5))
add_textbox(slide, Inches(0.7), Inches(1.7), Inches(5.4), Inches(0.35),
            "❌ 传统方案：直接向量化对比", font_size=Pt(14), bold=True, color=C_RED)

# Traditional flow
trad_items = ["宕机日志", "BGE-M3\n编码", "Cosine\nSimilarity\n≈ 0.3", "BGE-M3\n编码", "Commit\n描述"]
trad_x = Inches(0.7)
for i, item in enumerate(trad_items):
    add_rounded_rect(slide, trad_x, Inches(2.3), Inches(0.95), Inches(1.0), C_WHITE)
    add_textbox(slide, trad_x, Inches(2.4), Inches(0.95), Inches(0.8),
                item, font_size=Pt(9), bold=True, color=C_RED, alignment=PP_ALIGN.CENTER)
    if i < len(trad_items) - 1:
        add_textbox(slide, trad_x + Inches(0.95), Inches(2.6), Inches(0.25), Inches(0.3),
                    "→", font_size=Pt(14), color=C_RED, alignment=PP_ALIGN.CENTER)
    trad_x += Inches(1.1)

add_textbox(slide, Inches(0.7), Inches(3.5), Inches(5.4), Inches(0.7),
            '问题：宕机日志呈现「运行时现象」（如 NULL pointer at 0x28），\n'
            '补丁描述"代码修复逻辑"（如 add NULL check in slub allocator），\n'
            '两者在向量空间中距离远，直接比对相似度极低（~0.3）。',
            font_size=Pt(9.5), color=C_GRAY)

# Solution side (Right) - Green tint
add_rounded_rect(slide, Inches(6.8), Inches(1.6), Inches(6.0), Inches(2.8),
                 RGBColor(0xF0, 0xFF, 0xF4))
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(0.35),
            "✅ 本方案：对称根因嵌入（结构对齐）", font_size=Pt(14), bold=True, color=C_GREEN)

# Solution flow
sol_items = ["宕机日志", "RootCause\nAnalyzer\n→根因文本", "Cosine\nSimilarity\n≈ 0.85", "CommitRoot\nCauseBuilder\n→修复文本", "Commit\n信息"]
sol_x = Inches(7.0)
for i, item in enumerate(sol_items):
    add_rounded_rect(slide, sol_x, Inches(2.3), Inches(1.05), Inches(1.0), C_WHITE)
    add_textbox(slide, sol_x, Inches(2.4), Inches(1.05), Inches(0.8),
                item, font_size=Pt(9), bold=True, color=C_GREEN, alignment=PP_ALIGN.CENTER)
    if i < len(sol_items) - 1:
        if i == 1:
            add_textbox(slide, sol_x + Inches(1.05), Inches(2.55), Inches(0.25), Inches(0.4),
                        "≈0.85", font_size=Pt(12), bold=True, color=C_GREEN, alignment=PP_ALIGN.CENTER)
        else:
            add_textbox(slide, sol_x + Inches(1.05), Inches(2.6), Inches(0.25), Inches(0.3),
                        "→", font_size=Pt(14), color=C_GREEN, alignment=PP_ALIGN.CENTER)
    sol_x += Inches(1.1)

add_textbox(slide, Inches(7.0), Inches(3.5), Inches(5.6), Inches(0.7),
            '核心思路：将「现象/描述」统一转化为「根因/修复」结构对称文本，\n'
            '嵌入文本格式对齐：BugType + Subsystem + RootCause + FixPattern\n'
            '同语义空间编码，两端精准对齐，相似度从 0.3 → 0.85。',
            font_size=Pt(9.5), color=C_GRAY)

# Bottom: Embedding text structure example
add_rounded_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), C_WHITE)
add_textbox(slide, Inches(0.7), Inches(4.8), Inches(11.9), Inches(0.35),
            "▎嵌入文本结构对称示例（离线 / 在线两端统一格式）", font_size=Pt(13), bold=True, color=C_DARK_BG)

# Offline embedding text example
add_rounded_rect(slide, Inches(0.7), Inches(5.25), Inches(5.8), Inches(1.75),
                 RGBColor(0xEB, 0xF4, 0xFF))
add_textbox(slide, Inches(0.9), Inches(5.3), Inches(5.4), Inches(0.3),
            "离线侧（Commit 索引）— CommitRootCauseBuilder 生成", font_size=Pt(9), bold=True, color=C_BLUE)
add_textbox(slide, Inches(0.9), Inches(5.7), Inches(5.4), Inches(1.2),
            "BugType: null_pointer_dereference\n"
            "Subsystem: mm/slub\n"
            "RootCause: Missing NULL check before pointer dereference in slub allocator\n"
            "FixPattern: Add NULL pointer validation before accessing struct member\n"
            "Symptoms: NULL pointer dereference, kernel panic, unable to handle",
            font_size=Pt(9), color=C_BLUE, font_name="Courier New")

# Online embedding text example
add_rounded_rect(slide, Inches(6.8), Inches(5.25), Inches(5.8), Inches(1.75),
                 RGBColor(0xFF, 0xF4, 0xE6))
add_textbox(slide, Inches(7.0), Inches(5.3), Inches(5.4), Inches(0.3),
            "在线侧（宕机诊断）— RootCauseAnalyzer 生成", font_size=Pt(9), bold=True, color=C_ACCENT)
add_textbox(slide, Inches(7.0), Inches(5.7), Inches(5.4), Inches(1.2),
            "BugType: null_pointer_dereference\n"
            "Subsystem: mm\n"
            "RootCause: NULL pointer dereference without validity check\n"
            "FixPattern: Add NULL pointer validation\n"
            "Symptoms: BUG: unable to handle kernel NULL pointer at 0x28, Call Trace: slub_free",
            font_size=Pt(9), color=C_ACCENT, font_name="Courier New")

add_textbox(slide, Inches(3.5), Inches(5.25), Inches(6.3), Inches(0.25),
            "⬆ 两端输出结构完全对称，在相同语义空间编码 ⬆",
            font_size=Pt(9), bold=True, color=C_GREEN, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 5)

# ============================================================
# SLIDE 6: SYMMETRIC EMBEDDING ARCHITECTURE DETAIL
# ============================================================
print("Building Slide 6 - Symmetric Embedding Architecture...")
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_section_header(slide, 3, "对称嵌入双端设计：离线轻量 vs 在线精准",
                   "Dual-End Design: Lightweight Offline vs. Accurate Online")

# Comparison table
comp_headers = ["对比维度", "离线侧（Commit 索引）", "在线侧（宕机诊断）"]
comp_cols = [Inches(1.8), Inches(5.0), Inches(5.0)]
comp_rows = [
    ["核心模块", "CommitRootCauseBuilder", "RootCauseAnalyzer"],
    ["分析层次", "3层（模板查表→Diff规则→置信度）", "4层（专家规则→调用栈→Bug抽象→兜底）"],
    ["处理耗时", "3-5 ms/commit（百万级吞吐）", "~100 ms/次（在线实时）"],
    ["输入数据", "CommitInfo（title/body/diff/files）", "CrashFeature（panic_type/call_trace/subsystem）"],
    ["规则数量", "25条 DIFF_RULES（锁/RCU/内存等）", "28条 EXPERT_RULES + 10种Bug模式"],
    ["核心能力", "轻量快速，支持145万+commit索引", "深度精确，4层推断+LLM协同推理"],
    ["输出格式", "embedding_text → BGE-M3 → Milvus", "retrieval_query → BGE-M3 → 向量检索"],
]

add_table(slide, Inches(0.5), Inches(1.6), comp_cols, comp_headers, comp_rows, C_BLUE)

# Bottom: Visual showing symmetry
add_rounded_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.5), C_WHITE)
add_textbox(slide, Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.35),
            "▎对称嵌入文本结构模板（突出「结构对称」）", font_size=Pt(13), bold=True, color=C_DARK_BG)

# Show template structure
template_fields = [
    ("BugType", "bug_type 枚举值（21种）", C_RED),
    ("Subsystem", "subsystem 枚举值（28个）", C_BLUE),
    ("RootCause", "根因描述（规则推断+LLM生成）", C_ACCENT),
    ("FixPattern", "修复模式（锁添加/空指针检查/RCU保护...）", C_GREEN),
    ("Symptoms", "关键症状关键词（Panic/Error/Call Trace...）", C_GRAY),
]

field_x = Inches(0.9)
for field_name, field_desc, field_color in template_fields:
    # Field box
    add_rounded_rect(slide, field_x, Inches(5.45), Inches(2.2), Inches(1.6), C_WHITE)
    add_line(slide, field_x, Inches(5.45), Inches(2.2), Pt(3), field_color)
    add_textbox(slide, field_x, Inches(5.55), Inches(2.2), Inches(0.3),
                field_name, font_size=Pt(12), bold=True, color=field_color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, field_x, Inches(5.95), Inches(2.2), Inches(0.8),
                field_desc, font_size=Pt(9), color=C_GRAY, alignment=PP_ALIGN.CENTER)
    field_x += Inches(2.45)

add_textbox(slide, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.25),
            "★ 核心关键：离线侧与在线侧使用完全相同的字段结构生成 embedding 文本，确保 BGE-M3 编码后在相同语义空间可比较 ★",
            font_size=Pt(9.5), bold=True, color=C_ACCENT, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 6)

# ============================================================
# SLIDE 7: OFFLINE PHASE - COMMIT KNOWLEDGE BASE
# ============================================================
print("Building Slide 7 - Offline Phase...")
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_section_header(slide, 4, "离线核心：百万级 Commit 的结构化与向量化",
                   "Offline Core: Structured Indexing of Million-Scale Commits")

# Main flow diagram
flow_modules = [
    ("Linux Kernel\nGit Repo\n(145万+ commits)", C_DARK_BG),
    ("PyDriller\n流式遍历\nO(1)内存", C_BLUE),
    ("结构化提取\nSubsystem(28)\nBugType(21)", C_BLUE),
    ("CommitRootCause\nBuilder\n3层分析引擎", C_ACCENT),
    ("BGE-M3\n向量编码\nGPU加速", C_ACCENT),
    ("Milvus/FAISS\n向量存储\nIP度量索引", C_GREEN),
]

fx = Inches(0.3)
for i, (name, color) in enumerate(flow_modules):
    box = add_rounded_rect(slide, fx, Inches(1.6), Inches(2.0), Inches(1.8), C_WHITE)
    box.line.color.rgb = color
    box.line.width = Pt(2)
    add_textbox(slide, fx, Inches(1.75), Inches(2.0), Inches(1.5),
                name, font_size=Pt(10), bold=True, color=color, alignment=PP_ALIGN.CENTER)

    if i < len(flow_modules) - 1:
        add_textbox(slide, fx + Inches(2.0), Inches(2.2), Inches(0.2), Inches(0.3),
                    "▸", font_size=Pt(20), bold=True, color=C_ACCENT, alignment=PP_ALIGN.CENTER)

    fx += Inches(2.2)

# Key rules detail
add_rounded_rect(slide, Inches(0.3), Inches(3.7), Inches(6.2), Inches(3.5), C_WHITE)
add_textbox(slide, Inches(0.5), Inches(3.8), Inches(5.8), Inches(0.35),
            "▎CommitRootCauseBuilder 轻量引擎（3层，3-5ms/commit）", font_size=Pt(13), bold=True, color=C_BLUE)

add_textbox(slide, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.7),
            "第1层：BUG_TEMPLATE 查表匹配\n"
            "  • 10种预定义Bug模式模板（NULL pointer / UAF / deadlock / race...）\n"
            "  • 正则匹配 Commit message body + diff 特征代码\n\n"
            "第2层：DIFF_RULES 规则分析（25条）\n"
            "  • 锁相关规则（8条）：spin_lock_irqsave检测、mutex_lock模式匹配\n"
            "  • RCU规则（4条）：rcu_read_lock配对检查、rcu_dereference保护\n"
            "  • 引用计数规则（4条）：kref_get/put配对、refcount_t溢出检查\n"
            "  • 内存规则（5条）：kmalloc/kfree配对、slab分配器使用模式\n"
            "  • 其他规则（4条）：错误路径处理、空指针检查、边界检查\n\n"
            "第3层：置信度评估\n"
            "  • 综合规则命中数 + Fixes:/Cc:stable标签 → 置信度 0~1.0",
            font_size=Pt(9), color=C_GRAY)

# Performance indicators
add_rounded_rect(slide, Inches(6.8), Inches(3.7), Inches(6.0), Inches(3.5), C_WHITE)
add_textbox(slide, Inches(7.0), Inches(3.8), Inches(5.6), Inches(0.35),
            "▎关键性能指标", font_size=Pt(13), bold=True, color=C_ACCENT)

# Metric cards
metrics = [
    ("3-5ms", "单条Commit\n处理耗时", C_BLUE),
    ("145万+", "全量Commit\n索引规模", C_ACCENT),
    ("O(1)", "内存占用\n流式生成器", C_GREEN),
    ("~3.1GB", "全量FAISS\n索引大小", C_BLUE),
]

mx = Inches(7.0)
for val, label, color in metrics:
    add_rounded_rect(slide, mx, Inches(4.3), Inches(1.35), Inches(1.5), C_LIGHT_BG)
    add_textbox(slide, mx, Inches(4.45), Inches(1.35), Inches(0.5),
                val, font_size=Pt(22), bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, mx, Inches(5.0), Inches(1.35), Inches(0.6),
                label, font_size=Pt(9), color=C_GRAY, alignment=PP_ALIGN.CENTER)
    mx += Inches(1.45)

add_textbox(slide, Inches(7.0), Inches(5.95), Inches(5.6), Inches(1.0),
            "其他特性：\n"
            "• 断点续跑：--resume 参数支持中断恢复\n"
            "• 批量GPU编码：batch_size可配置\n"
            "• 双后端自动降级：Milvus → FAISS\n"
            "• 版本标注：commit date → kernel release version",
            font_size=Pt(9), color=C_GRAY)

add_page_number(slide, 7)

# ============================================================
# SLIDE 8: ONLINE PHASE - FOUR-STAGE RETRIEVAL
# ============================================================
print("Building Slide 8 - Online Four-Stage Retrieval...")
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_section_header(slide, 5, "在线核心：秒级精准检索的四阶段链路",
                   "Online Core: Four-Stage Precision Retrieval Pipeline")

# Funnel visualization using shapes
funnel_stages = [
    ("Stage 1: Recall\n向量快速召回", "Top-100", "< 50ms", C_BLUE, Inches(2.5)),
    ("Stage 2: Filter\n规则硬过滤", "50-80", "+10ms", C_GREEN, Inches(3.8)),
    ("Stage 3: Rerank\n深度语义重排", "Top-20", "+500ms", C_ACCENT, Inches(5.1)),
    ("Stage 4: LLM Judge\n因果推理评分", "Top-3~5", "+2-10s", C_RED, Inches(6.4)),
]

# Draw funnel stages as trapezoid-like shapes
prev_w = Inches(10.0)
for i, (name, candidates, latency, color, y) in enumerate(funnel_stages):
    w = prev_w - Inches(0.8 * i)
    x = (SLIDE_W - w) / 2

    # Stage box
    box = add_rounded_rect(slide, x, y, w, Inches(0.85), C_WHITE)
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Left accent
    add_line(slide, x, y, Inches(0.06), Inches(0.85), color)

    # Stage name (left)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.5), Inches(0.75),
                name, font_size=Pt(11), bold=True, color=color)

    # Candidate count (center-right)
    add_textbox(slide, x + w - Inches(3.0), y + Inches(0.05), Inches(1.2), Inches(0.75),
                candidates, font_size=Pt(18), bold=True, color=color,
                alignment=PP_ALIGN.CENTER)

    # Latency (right)
    add_textbox(slide, x + w - Inches(1.7), y + Inches(0.05), Inches(1.4), Inches(0.75),
                latency, font_size=Pt(10), color=C_GRAY, alignment=PP_ALIGN.CENTER)

    prev_w = w

# Upward arrow between stages
for i in range(3):
    mid_y = Inches(2.5) + i * Inches(1.3) + Inches(0.85)
    add_textbox(slide, Inches(6.4), mid_y, Inches(0.5), Inches(0.4),
                "▼", font_size=Pt(16), color=C_ACCENT, alignment=PP_ALIGN.CENTER)

# Bottom detail boxes
details = [
    ("输入处理", "dmesg文本 → 20+Panic模式正则\nvmcore → drgn内核对象提取\n自然语言 → LLM特征提取\n→ CrashFeature", C_BLUE),
    ("检索模式", "fast: <100ms (仅Recall)\nstandard: <1s (Recall+Rerank)\ndeep: 2-10s (全四阶段)\n灵活按需选择", C_GREEN),
    ("输出报告", "Markdown/JSON格式\n补丁列表 + 匹配理由\n置信度评分 + 因果链\n可解释诊断报告", C_ACCENT),
]

dx = Inches(0.5)
for title, content, color in details:
    add_rounded_rect(slide, dx, Inches(7.4), Inches(4.0), Inches(-1.6), None)  # placeholder

    add_rounded_rect(slide, dx, Inches(1.45), Inches(3.9), Inches(1.8), C_WHITE)
    add_line(slide, dx + Inches(0.02), Inches(1.5), Inches(0.05), Inches(1.7), color)
    add_textbox(slide, dx + Inches(0.2), Inches(1.55), Inches(3.5), Inches(0.3),
                title, font_size=Pt(12), bold=True, color=color)
    add_textbox(slide, dx + Inches(0.2), Inches(1.9), Inches(3.5), Inches(1.2),
                content, font_size=Pt(9.5), color=C_GRAY)
    dx += Inches(4.2)

add_page_number(slide, 8)

# Wait - fix the bottom boxes position
# They should be at the bottom, not overlapping the funnel
# Let me fix this in the next save

# ============================================================
# SLIDE 9: TECHNICAL HIGHLIGHTS & PERFORMANCE
# ============================================================
print("Building Slide 9 - Technical Highlights...")
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_section_header(slide, 6, "核心技术亮点与关键性能指标",
                   "Technical Highlights & Key Performance Metrics")

# Left: Technical Highlights
add_textbox(slide, Inches(0.5), Inches(1.55), Inches(6), Inches(0.35),
            "▎核心技术亮点", font_size=Pt(15), bold=True, color=C_DARK_BG)

highlights = [
    ("⭐ 语义对称嵌入", "解决核心语义鸿沟问题\n相似度从 0.3 → 0.85（提升 2.8x）", C_ACCENT),
    ("🔍 四阶段检索", "兼顾召回率（95%+）与效率（秒级响应）\n100→80→20→5 漏斗式精准收敛", C_BLUE),
    ("📋 版本感知", "内核版本敏感的过滤与匹配\n覆盖 kernel 4.9 ~ 6.13 全版本标注", C_GREEN),
    ("🧩 多模态融合", "drgn解析 vmcore + LLM分析 dmesg\n+ 规则匹配，异构数据统一处理", C_BLUE),
    ("🛡️ 容错降级", "Milvus→FAISS、LLM→规则引擎\n三级自动降级，核心功能不受影响", C_GREEN),
    ("🆓 免费可用", "集成Ollama本地大模型\n无需API Key即可使用LLM分析", C_ACCENT),
]

hy = Inches(2.0)
for title, desc, color in highlights:
    add_rounded_rect(slide, Inches(0.5), hy, Inches(5.8), Inches(0.75), C_WHITE)
    add_line(slide, Inches(0.52), hy + Inches(0.05), Inches(0.05), Inches(0.65), color)
    add_textbox(slide, Inches(0.75), hy + Inches(0.05), Inches(5.3), Inches(0.3),
                title, font_size=Pt(11.5), bold=True, color=color)
    add_textbox(slide, Inches(0.75), hy + Inches(0.35), Inches(5.3), Inches(0.35),
                desc, font_size=Pt(9), color=C_GRAY)
    hy += Inches(0.85)

# Right: Performance Metrics
add_textbox(slide, Inches(6.8), Inches(1.55), Inches(6), Inches(0.35),
            "▎关键性能指标", font_size=Pt(15), bold=True, color=C_DARK_BG)

perf_headers = ["指标", "目标值", "说明"]
perf_cols = [Inches(1.8), Inches(1.5), Inches(2.5)]
perf_rows = [
    ["离线索引耗时", "3-5ms/条", "百万级全量索引可行"],
    ["在线诊断耗时", "~100ms/次", "根因分析+检索+LLM判优"],
    ["检索召回率", "≥95% (Top-100)", "覆盖95%以上相关补丁"],
    ["匹配精度", "≥80% (Top-5)", "Top-5命中修复补丁比例"],
    ["端到端延迟", "<3s (standard)", "提交日志到返回结果"],
    ["API响应", "<500ms", "统计/搜索接口响应时间"],
    ["并发支持", "50+ QPS", "FastAPI异步处理能力"],
    ["索引规模", "312,632条", "全量模式覆盖4.9~6.13"],
]

add_table(slide, Inches(6.8), Inches(2.0), perf_cols, perf_headers, perf_rows, C_ACCENT)

# Bottom mode comparison
add_rounded_rect(slide, Inches(6.8), Inches(5.6), Inches(5.8), Inches(1.7), C_WHITE)
add_textbox(slide, Inches(7.0), Inches(5.7), Inches(5.4), Inches(0.3),
            "▎三种检索模式对比", font_size=Pt(12), bold=True, color=C_DARK_BG)

mode_headers = ["模式", "Recall", "Filter", "Rerank", "LLM", "延迟"]
mode_cols = [Inches(0.9)] + [Inches(0.7)] * 4 + [Inches(1.0)]
mode_rows = [
    ["fast", "✅ Top-50", "✅", "❌", "❌", "< 100ms"],
    ["standard", "✅ Top-100", "✅", "✅ BGE", "❌", "< 1s"],
    ["deep", "✅ Top-200", "✅", "✅ BGE", "✅ LLM", "2-10s"],
]

add_table(slide, Inches(7.0), Inches(6.15), mode_cols, mode_headers, mode_rows, C_BLUE)

add_page_number(slide, 9)

# ============================================================
# SLIDE 10: SUMMARY & APPLICATION VALUE
# ============================================================
print("Building Slide 10 - Summary & Value...")
slide = add_blank_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_LIGHT_BG)
add_section_header(slide, 7, "方案总结与产业应用价值",
                   "Solution Summary & Industrial Application Value")

# Left: Solution Summary
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(4.8), C_WHITE)
add_textbox(slide, Inches(0.7), Inches(1.6), Inches(5.4), Inches(0.35),
            "▎方案总结", font_size=Pt(16), bold=True, color=C_DARK_BG)

add_textbox(slide, Inches(0.7), Inches(2.1), Inches(5.4), Inches(4.0),
            "核心解决：\n"
            "  Linux内核宕机补丁匹配的 「效率低、精度差、依赖专家」三大问题\n\n"
            "技术路径：\n"
            "  RAG + 内核领域知识 + LLM 深度结合\n"
            "  离线+在线双阶段架构，对称嵌入+四阶段检索\n\n"
            "核心优势：\n"
            "  ✅ 高精度：Top-5 命中精度 80%+，语义相似度提升 2.8x\n"
            "  ✅ 高效率：秒级端到端响应（<3s），全自动无需人工\n"
            "  ✅ 可解释：LLM 生成推荐理由 + 因果推理链\n"
            "  ✅ 可扩展：适配多版本内核（4.9~6.13），支持新增Bug类型\n"
            "  ✅ 低成本：支持Ollama本地免费模型，Docker一键部署\n\n"
            "关键创新：\n"
            "  ⭐ 对称 Root Cause Embedding（核心突破）\n"
            "  ⭐ 四阶段级联检索架构\n"
            "  ⭐ 28条专家规则 + 10种Bug模式知识图谱",
            font_size=Pt(10.5), color=C_TEXT)

# Right: Application Value
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.8), C_WHITE)
add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.35),
            "▎应用价值与前景", font_size=Pt(16), bold=True, color=C_DARK_BG)

# Value cards
values = [
    ("🚀 运维提效", "故障分析从 小时级→分钟级\n降低对资深内核专家的依赖\n自动化诊断，7×24可用", C_GREEN),
    ("💰 成本节约", "减少数据中心宕机带来的\n业务中断损失（MTTR↓）\n降低运维人力成本", C_ACCENT),
    ("🔧 可扩展性", "适配Linux 4.9~6.13全版本\n支持新增子系统/Bug类型\n知识库持续进化", C_BLUE),
    ("🏭 产业应用", "云计算/数据中心内核运维\n边缘计算节点自动诊断\nLinux发行版Bug分类", C_BLUE),
]

vy = Inches(2.2)
for title, desc, color in values:
    add_rounded_rect(slide, Inches(7.0), vy, Inches(5.6), Inches(0.9), C_LIGHT_BG)
    add_line(slide, Inches(7.02), vy + Inches(0.05), Inches(0.05), Inches(0.8), color)
    add_textbox(slide, Inches(7.2), vy + Inches(0.05), Inches(5.2), Inches(0.3),
                title, font_size=Pt(12), bold=True, color=color)
    add_textbox(slide, Inches(7.2), vy + Inches(0.35), Inches(5.2), Inches(0.5),
                desc, font_size=Pt(9.5), color=C_GRAY)
    vy += Inches(1.05)

# Bottom banner
add_rounded_rect(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6), C_ACCENT)
add_textbox(slide, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.5),
            '🎯 技术赋能内核运维智能化 — 从「人工经验驱动」到「数据+AI驱动」的范式升级',
            font_size=Pt(14), bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 10)

# ============================================================
# SLIDE 11: Q&A
# ============================================================
print("Building Slide 11 - Q&A...")
slide = add_blank_slide()

# Dark background
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), C_ACCENT)
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), C_ACCENT)

# Linux penguin icon
add_textbox(slide, Inches(5.5), Inches(1.0), Inches(2.3), Inches(1.5),
            "🐧", font_size=Pt(60), alignment=PP_ALIGN.CENTER, color=C_WHITE)

# Main text
add_textbox(slide, Inches(2.0), Inches(2.5), Inches(9.3), Inches(1.2),
            "感谢聆听", font_name=F_TITLE, font_size=Pt(52), bold=True, color=C_WHITE,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2.0), Inches(3.7), Inches(9.3), Inches(0.6),
            "Q & A", font_name=F_EN, font_size=Pt(36), bold=True, color=C_ACCENT,
            alignment=PP_ALIGN.CENTER)

# Project name
add_textbox(slide, Inches(2.0), Inches(4.6), Inches(9.3), Inches(0.5),
            "Linux 内核宕机自动诊断与补丁匹配系统",
            font_size=Pt(18), color=RGBColor(0xA0, 0xAE, 0xC0), alignment=PP_ALIGN.CENTER)

# Bottom info
add_multiline_textbox(slide, Inches(3.0), Inches(5.8), Inches(7.3), Inches(1.0), [
    {'text': '项目团队：XXX    |    指导教师：XXX    |    2026年6月',
     'font_size': Pt(12), 'color': C_GRAY, 'alignment': PP_ALIGN.CENTER},
    {'text': '基于 RAG + 内核领域知识 + LLM 的自动化补丁匹配方案',
     'font_size': Pt(11), 'color': RGBColor(0x60, 0x6E, 0x80), 'alignment': PP_ALIGN.CENTER},
])

# Decorative lines
add_line(slide, Inches(5.0), Inches(4.2), Inches(3.3), Pt(2), C_ACCENT)
add_line(slide, Inches(5.5), Inches(5.5), Inches(2.3), Pt(1), RGBColor(0x40, 0x4E, 0x60))

add_page_number(slide, 11)

# ============================================================
# FIX SLIDE 8: Move bottom detail boxes down
# ============================================================
print("Fixing Slide 8 layout...")
slide8 = prs.slides[7]  # Slide 8 (0-indexed)

# Remove the overlapping detail boxes and recreate them at the bottom
# First, let's identify shapes that need moving
# The detail boxes were added at y=1.45 but funnel is at y=2.5-7.25
# They need to be at the very bottom

# Add bottom detail boxes using fresh coordinates
detail_data = [
    ("输入处理", "dmesg文本 → 20+Panic模式正则 | vmcore → drgn内核对象提取\n自然语言 → LLM特征提取 → 统一 CrashFeature 结构化表示", C_BLUE),
    ("检索模式", "fast: <100ms（仅Recall Top-50）\nstandard: <1s（Recall+Rerank）| deep: 2-10s（全四阶段+LLM）\n按需灵活选择，兼顾速度与精度", C_GREEN),
    ("输出报告", "Markdown / JSON 格式报告\n补丁列表 + 匹配理由 + 置信度评分 + 因果推理链\n可解释诊断报告，助力快速定位修复方案", C_ACCENT),
]

dx = Inches(0.5)
for title, content, color in detail_data:
    add_rounded_rect(slide8, dx, Inches(5.9), Inches(3.9), Inches(1.2), C_WHITE)
    add_line(slide8, dx + Inches(0.02), Inches(5.95), Inches(0.05), Inches(1.1), color)
    add_textbox(slide8, dx + Inches(0.2), Inches(5.95), Inches(3.5), Inches(0.25),
                title, font_size=Pt(11), bold=True, color=color)
    add_textbox(slide8, dx + Inches(0.2), Inches(6.2), Inches(3.5), Inches(0.8),
                content, font_size=Pt(8.5), color=C_GRAY)
    dx += Inches(4.2)

# ============================================================
# SAVE
# ============================================================
output_path = '/home/lkx/桌面/Linux内核宕机自动诊断与补丁匹配系统_比赛汇报.pptx'
prs.save(output_path)
print(f"\n✅ PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
